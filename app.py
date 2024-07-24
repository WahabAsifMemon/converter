import os
import logging
from flask import Flask, request, send_file, render_template, jsonify
from pdf2image import convert_from_path
import zipfile
import io
import datetime
from pdf2docx import Converter
import tabula
import pandas as pd
import fitz  # PyMuPDF library for PDF to PDF/A conversion
from docx2pdf import convert as docx2pdf_convert
from werkzeug.utils import secure_filename
import subprocess
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PyPDF2 import PdfMerger, PdfReader, PdfWriter, PdfFileReader, PdfFileWriter
from fpdf import FPDF
from PIL import Image
import PyPDF2

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'output'

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

if not os.path.exists(OUTPUT_FOLDER):
    os.makedirs(OUTPUT_FOLDER)

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s %(message)s')


def clear_folder(folder_path):
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        try:
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
            elif os.path.isdir(file_path):
                os.rmdir(file_path)
        except Exception as e:
            logging.error(f'Failed to delete {file_path}. Reason: {e}')


@app.route('/')
def home():
    return render_template('index.html')

@app.route('/pdf-to-jpg')
def pdf_to_jpg():
    return render_template('pdf-to-jpg.html')

@app.route('/pdf-to-word')
def pdf_to_word():
    return render_template('pdf-to-word.html')

@app.route('/pdf-to-excel')
def pdf_to_excel():
    return render_template('pdf-to-excel.html')

@app.route('/pdf-to-ppt')
def pdf_to_ppt():
    return render_template('pdf-to-ppt.html')

@app.route('/pdf-to-pdfa')
def pdf_to_pdfa():
    return render_template('pdf-to-pdfa.html')


@app.route('/html-to-pdf')
def html_to_pdf():
    return render_template('html-to-pdf.html')



@app.route('/ppt-to-pdf')
def ppt_to_pdf():
    return render_template('ppt-to-pdf.html')

@app.route('/excel-to-pdf')
def excel_to_pdf():
    return render_template('excel-to-pdf.html')

@app.route('/jpg-to-pdf')
def jpg_to_pdf():
    return render_template('jpg-to-pdf.html')

@app.route('/word-to-pdf')
def word_to_pdf():
    return render_template('word-to-pdf.html')

@app.route('/merge-pdf')
def merge_pdf():
    return render_template('merge-pdf.html')

@app.route('/compress-pdf')
def compress_pdf():
    return render_template('compress-pdf.html')

@app.route('/split-pdf')
def split_pdf():
    return render_template('split-pdf.html')


@app.route('/unlock-pdf')
def unlock_pdf():
    return render_template('unlock-pdf.html')


ALLOWED_EXTENSIONS = {'pdf'}
ALLOWED_EXTENSIONS_DOCX = {'docx'}
ALLOWED_EXTENSIONS_PPTX = {'pptx'}
ALLOWED_EXTENSIONS_IMG = {'jpg', 'jpeg'}



def allowed_file(filename, allowed_extensions):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions


def pdf_allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# START PDF TO JPG
@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            logging.error('No file part in the request')
            return jsonify({'error': 'No file part'}), 400

        file = request.files['file']

        if file.filename == '':
            logging.error('No selected file')
            return jsonify({'error': 'No selected file'}), 400

        if file and pdf_allowed_file(file.filename):
            clear_folder(UPLOAD_FOLDER)
            clear_folder(OUTPUT_FOLDER)

            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            file.save(file_path)
            images = convert_from_path(file_path)
            image_files = []

            for i, image in enumerate(images):
                image_filename = f'page_{i + 1}.jpg'
                image_path = os.path.join(OUTPUT_FOLDER, image_filename)
                image.save(image_path, 'JPEG')
                image_files.append(image_filename)

            original_filename = os.path.splitext(file.filename)[0]

            return jsonify({'images': image_files, 'original_filename': original_filename})

        else:
            logging.error('Invalid file type, only PDF files are allowed')
            return jsonify({'error': 'Invalid file type, only PDF files are allowed'}), 400

    except Exception as e:
        logging.error(f'Error during file upload: {e}')
        return jsonify({'error': f'File upload failed: {str(e)}'}), 500
# END PDF TO JPG

# START PDF TO WORD
@app.route('/upload-pdf-to-word', methods=['POST'])
def upload_pdf_to_word():
    try:
        if 'file' not in request.files:
            logging.error('No file part in the request')
            return jsonify({'error': 'No file part'}), 400

        file = request.files['file']

        if file.filename == '':
            logging.error('No selected file')
            return jsonify({'error': 'No selected file'}), 400

        if file and allowed_file(file.filename):
            clear_folder(UPLOAD_FOLDER)
            clear_folder(OUTPUT_FOLDER)

            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            file.save(file_path)

            # Convert PDF to Word
            docx_file_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(file.filename)[0]}.docx")
            cv = Converter(file_path)
            cv.convert(docx_file_path)
            cv.close()

            return jsonify({'filename': f"{os.path.splitext(file.filename)[0]}.docx"}), 200

        else:
            logging.error('Invalid file type, only PDF files are allowed')
            return jsonify({'error': 'Invalid file type, only PDF files are allowed'}), 400

    except Exception as e:
        logging.error(f'Error during file upload: {e}')
        return jsonify({'error': f'File upload failed: {str(e)}'}), 500
# END PDF TO WORD


# START PDF TO EXCEL
@app.route('/upload-pdf-to-excel', methods=['POST'])
def upload_pdf_to_excel():
    try:
        if 'file' not in request.files:
            logging.error('No file part in the request')
            return jsonify({'error': 'No file part'}), 400

        file = request.files['file']

        if file.filename == '':
            logging.error('No selected file')
            return jsonify({'error': 'No selected file'}), 400

        if file and allowed_file(file.filename):
            clear_folder(UPLOAD_FOLDER)
            clear_folder(OUTPUT_FOLDER)

            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            file.save(file_path)

            # Convert PDF to Excel
            df_list = tabula.read_pdf(file_path, pages='all', multiple_tables=True)
            excel_file_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(file.filename)[0]}.xlsx")

            with pd.ExcelWriter(excel_file_path) as writer:
                for i, df in enumerate(df_list):
                    df.to_excel(writer, sheet_name=f'Sheet{i+1}', index=False)

            return jsonify({'filename': f"{os.path.splitext(file.filename)[0]}.xlsx"}), 200

        else:
            logging.error('Invalid file type, only PDF files are allowed')
            return jsonify({'error': 'Invalid file type, only PDF files are allowed'}), 400

    except Exception as e:
        logging.error(f'Error during file upload: {e}')
        return jsonify({'error': f'File upload failed: {str(e)}'}), 500
# END PDF TO EXCEL


# START PDF TO PDF/A
@app.route('/upload-pdf-to-pdfa', methods=['POST'])
def upload_pdf_to_pdfa():
    try:
        if 'file' not in request.files:
            logging.error('No file part in the request')
            return jsonify({'error': 'No file part'}), 400

        file = request.files['file']

        if file.filename == '':
            logging.error('No selected file')
            return jsonify({'error': 'No selected file'}), 400

        if file and allowed_file(file.filename):
            clear_folder(UPLOAD_FOLDER)
            clear_folder(OUTPUT_FOLDER)

            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            file.save(file_path)

            # Convert PDF to PDF/A
            pdfa_file_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(file.filename)[0]}_pdfa.pdf")
            convert_to_pdfa(file_path, pdfa_file_path)

            return jsonify({'filename': f"{os.path.splitext(file.filename)[0]}_pdfa.pdf"}), 200

        else:
            logging.error('Invalid file type, only PDF files are allowed')
            return jsonify({'error': 'Invalid file type, only PDF files are allowed'}), 400

    except Exception as e:
        logging.error(f'Error during file upload: {e}')
        return jsonify({'error': f'File upload failed: {str(e)}'}), 500
# END PDF TO PDF/A


@app.route('/upload-pdf-to-ppt', methods=['POST'])
def upload_pdf_to_ppt():
    try:
        if 'file' not in request.files:
            logging.error('No file part in the request')
            return jsonify({'error': 'No file part'}), 400

        file = request.files['file']

        if file.filename == '':
            logging.error('No selected file')
            return jsonify({'error': 'No selected file'}), 400

        if file and allowed_file(file.filename, ALLOWED_EXTENSIONS):
            clear_folder(UPLOAD_FOLDER)
            clear_folder(OUTPUT_FOLDER)

            file_path = os.path.join(UPLOAD_FOLDER, secure_filename(file.filename))
            file.save(file_path)

            # Convert PDF to images (one image per page)
            images = convert_from_path(file_path, 300)  # 300 DPI for good quality

            # Create a PowerPoint presentation
            ppt_file_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(file.filename)[0]}.pptx")
            prs = Presentation()

            for i, image in enumerate(images):
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
                slide.shapes.add_picture(io.BytesIO(image._repr_png_()), Inches(0), Inches(0), width=Inches(10))

            prs.save(ppt_file_path)

            return jsonify({'filename': f"{os.path.splitext(file.filename)[0]}.pptx"}), 200

        else:
            logging.error('Invalid file type, only PDF or PPTX files are allowed')
            return jsonify({'error': 'Invalid file type, only PDF or PPTX files are allowed'}), 400

    except Exception as e:
        logging.error(f'Error during file upload: {e}')
        return jsonify({'error': f'File upload failed: {str(e)}'}), 500


@app.route('/upload-ppt-to-pdf', methods=['POST'])
def upload_ppt_to_pdf():
    try:
        if 'file' not in request.files:
            logging.error('No file part in the request')
            return jsonify({'error': 'No file part'}), 400

        file = request.files['file']

        if file.filename == '':
            logging.error('No selected file')
            return jsonify({'error': 'No selected file'}), 400

        if file and allowed_file(file.filename, {'ppt', 'pptx'}):
            clear_folder(UPLOAD_FOLDER)
            clear_folder(OUTPUT_FOLDER)

            file_path = os.path.join(UPLOAD_FOLDER, secure_filename(file.filename))
            file.save(file_path)

            pdf_file_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(file.filename)[0]}.pdf")

            # Convert using LibreOffice
            cmd = f'libreoffice --headless --convert-to pdf --outdir {OUTPUT_FOLDER} {file_path}'
            subprocess.run(cmd, shell=True, check=True)

            return jsonify({'filename': f"{os.path.splitext(file.filename)[0]}.pdf"}), 200

        else:
            logging.error('Invalid file type, only PPT or PPTX files are allowed')
            return jsonify({'error': 'Invalid file type, only PPT or PPTX files are allowed'}), 400

    except subprocess.CalledProcessError as e:
        logging.error(f'LibreOffice conversion error: {e}')
        return jsonify({'error': f'File upload failed: {str(e)}'}), 500
    except Exception as e:
        logging.error(f'Error during file upload: {e}')
        return jsonify({'error': f'File upload failed: {str(e)}'}), 500


@app.route('/upload-excel-to-pdf', methods=['POST'])
def upload_excel_to_pdf():
    try:
        if 'file' not in request.files:
            logging.error('No file part in the request')
            return jsonify({'error': 'No file part'}), 400

        file = request.files['file']

        if file.filename == '':
            logging.error('No selected file')
            return jsonify({'error': 'No selected file'}), 400

        if file and allowed_file(file.filename, {'xlsx', 'xls'}):
            clear_folder(UPLOAD_FOLDER)
            clear_folder(OUTPUT_FOLDER)

            file_path = os.path.join(UPLOAD_FOLDER, secure_filename(file.filename))
            file.save(file_path)

            pdf_file_path = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(file.filename)[0]}.pdf")

            # Convert Excel to PDF
            convert_excel_to_pdf(file_path, pdf_file_path)

            return jsonify({'filename': f"{os.path.splitext(file.filename)[0]}.pdf"}), 200

        else:
            logging.error('Invalid file type, only Excel files are allowed')
            return jsonify({'error': 'Invalid file type, only Excel files are allowed'}), 400

    except Exception as e:
        logging.error(f'Error during file upload: {e}')
        return jsonify({'error': f'File upload failed: {str(e)}'}), 500

def convert_excel_to_pdf(excel_path, pdf_path):
    # Read the Excel file
    if excel_path.endswith('.xls'):
        df = pd.read_excel(excel_path, engine='xlrd')
    else:
        df = pd.read_excel(excel_path, engine='openpyxl')

    # Create a PDF document
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    # Define the row height and the starting position
    row_height = 20
    y = height - 40

    for i, row in df.iterrows():
        x = 40
        for value in row:
            c.drawString(x, y, str(value))
            x += 100  # Move to the next column position

        y -= row_height  # Move to the next row position

        # Check if we need to add a new page
        if y < row_height:
            c.showPage()
            y = height - 40

    c.save()


@app.route('/upload-and-merge', methods=['POST'])
def upload_and_merge():
    try:
        uploaded_files = request.files.getlist('file')

        if len(uploaded_files) < 2:
            return jsonify({'error': 'Please upload at least 2 PDF files to merge'}), 400

        clear_folder(UPLOAD_FOLDER)
        clear_folder(OUTPUT_FOLDER)

        file_paths = []
        for file in uploaded_files:
            if file and allowed_file(file.filename, {'pdf'}):
                filename = secure_filename(file.filename)
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                file.save(file_path)
                file_paths.append(file_path)
            else:
                logging.error(f'Invalid file type or no file selected: {file.filename}')
                return jsonify({'error': f'Invalid file type or no file selected: {file.filename}'}), 400

        merged_pdf_path = os.path.join(OUTPUT_FOLDER, 'merged_file.pdf')

        # Merge PDF files using PdfMerger and PdfReader
        merge_pdfs(file_paths, merged_pdf_path)

        return jsonify({'filename': 'merged_file.pdf'}), 200

    except Exception as e:
        logging.error(f'Error during file upload and merge: {e}')
        return jsonify({'error': f'File upload and merge failed: {str(e)}'}), 500

def merge_pdfs(input_paths, output_path):
    merger = PdfMerger()

    for path in input_paths:
        merger.append(PdfReader(open(path, 'rb')))

    merger.write(output_path)
    merger.close()


@app.route('/upload-and-compress', methods=['POST'])
def upload_and_compress():
    try:
        uploaded_file = request.files['file']

        if not uploaded_file:
            return jsonify({'error': 'No file uploaded'}), 400

        if not allowed_file(uploaded_file.filename, {'pdf'}):
            return jsonify({'error': 'Invalid file type. Please upload a PDF file'}), 400

        clear_folder(UPLOAD_FOLDER)
        clear_folder(OUTPUT_FOLDER)

        filename = secure_filename(uploaded_file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        uploaded_file.save(file_path)

        compressed_pdf_path = os.path.join(OUTPUT_FOLDER, 'compressed_file.pdf')

        # Compress PDF
        compress_pdf(file_path, compressed_pdf_path)

        return jsonify({'filename': 'compressed_file.pdf'}), 200

    except Exception as e:
        logging.error(f'Error during file upload and compression: {e}')
        return jsonify({'error': f'File upload and compression failed: {str(e)}'}), 500

def compress_pdf(input_path, output_path):
    with open(input_path, 'rb') as input_file:
        pdf_reader = PdfReader(input_file)
        pdf_writer = PdfWriter()

        for page_num in range(len(pdf_reader.pages)):
            pdf_writer.add_page(pdf_reader.pages[page_num])

        with open(output_path, 'wb') as output_file:
            pdf_writer.write(output_file)


@app.route('/upload-and-split', methods=['POST'])
def upload_and_split():
    try:
        uploaded_file = request.files['file']

        if not uploaded_file:
            return jsonify({'error': 'No file uploaded'}), 400

        if not allowed_file(uploaded_file.filename, {'pdf'}):
            return jsonify({'error': 'Invalid file type. Please upload a PDF file'}), 400

        clear_folder(UPLOAD_FOLDER)
        clear_folder(OUTPUT_FOLDER)

        filename = secure_filename(uploaded_file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        uploaded_file.save(file_path)

        # Split PDF
        split_pdf(file_path, OUTPUT_FOLDER)

        # Create a zip file of split PDFs
        zip_filename = 'split_files.zip'
        zip_path = os.path.join(OUTPUT_FOLDER, zip_filename)
        zipf = zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED)

        for root, _, files in os.walk(OUTPUT_FOLDER):
            for file in files:
                if file.endswith('.pdf'):
                    zipf.write(os.path.join(root, file), file)

        zipf.close()

        return jsonify({'zip_file': zip_filename}), 200

    except Exception as e:
        logging.error(f'Error during file upload and split: {e}')
        return jsonify({'error': f'File upload and split failed: {str(e)}'}), 500

def split_pdf(input_path, output_folder):
    with open(input_path, 'rb') as input_file:
        pdf_reader = PdfReader(input_file)
        num_pages = len(pdf_reader.pages)

        for page_num in range(num_pages):
            pdf_writer = PdfWriter()
            pdf_writer.add_page(pdf_reader.pages[page_num])

            output_path = os.path.join(output_folder, f'page_{page_num + 1}.pdf')

            with open(output_path, 'wb') as output_file:
                pdf_writer.write(output_file)

# CONVERTING PDF TO PDF/A
def convert_to_pdfa(input_path, output_path):
    # Using PyMuPDF to convert PDF to PDF/A
    try:
        doc = fitz.open(input_path)
        doc_pdf = fitz.open()
        doc_pdf.insert_pdf(doc)
        doc_pdf.save(output_path)
        doc_pdf.close()
    except Exception as e:
        logging.error(f'Error converting to PDF/A: {e}')
        raise
# CONVERTING PDF TO PDF/A

# START JPG TO PDF
@app.route('/upload-jpg-to-pdf', methods=['POST'])
def upload_jpg_to_pdf():
   print('Wanag')
# END JPG TO PDF



@app.route('/download_all')
def download_all():
    try:
        image_files = request.args.get('images').split(',')
        original_filename = request.args.get('filename', 'images')

        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for image in image_files:
                zip_file.write(os.path.join(OUTPUT_FOLDER, image), image)
        zip_buffer.seek(0)

        # Constructing the zip filename based on the original filename
        timestamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S')
        zip_filename = f'{original_filename}_{timestamp}.zip'

        return send_file(zip_buffer, mimetype='application/zip', as_attachment=True, download_name=zip_filename)
    except Exception as e:
        logging.error(f'Error during file download: {e}')
        return jsonify({'error': f'File download failed: {str(e)}'}), 500

@app.route('/download')
def download():
    filename = request.args.get('filename')
    if filename:
        file_path = os.path.join(OUTPUT_FOLDER, filename)
        return send_file(file_path, as_attachment=True)
    else:
        return jsonify({'error': 'Filename not provided'}), 400

def convert_jpg_to_pdf(input_folder, output_path):
    pdf = FPDF()
    image_list = []

    for root, _, files in os.walk(input_folder):
        for file in files:
            if file.endswith('.jpg') or file.endswith('.jpeg'):
                image_list.append(os.path.join(root, file))

    for image_path in sorted(image_list):
        image = Image.open(image_path)
        width, height = image.size
        pdf.add_page()
        pdf.image(image_path, 0, 0, width * 0.75, height * 0.75)

    pdf.output(output_path, 'F')

    @app.route('/upload-unlock-pdf', methods=['POST'])
    def upload_unlock_pdf():
        if 'file' not in request.files:
            return 'No file part'

        file = request.files['file']
        password = request.form.get('password')

        if file and allowed_file(file.filename) and password:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            # Unlock the PDF
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                if reader.is_encrypted:
                    try:
                        reader.decrypt(password)
                    except:
                        return 'Incorrect password or unable to decrypt PDF'

                    writer = PyPDF2.PdfWriter()
                    for page_num in range(len(reader.pages)):
                        writer.add_page(reader.pages[page_num])

                    unlocked_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], 'unlocked_' + filename)
                    with open(unlocked_pdf_path, 'wb') as out_f:
                        writer.write(out_f)

                    return send_file(unlocked_pdf_path, as_attachment=True, download_name='unlocked_' + filename)
                else:
                    return 'PDF is not encrypted'

        return 'Invalid file or missing password'


if __name__ == '__main__':
    app.run(debug=True)
